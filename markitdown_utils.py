"""
提供 MarkItDown 相關功能的實用函數
"""

import os
import tempfile
from pathlib import Path
import logging
from markitdown import MarkItDown
from openai import OpenAI, AuthenticationError
from typing import Dict, Any, Optional, List, Tuple
import subprocess
import sys
import base64
from io import BytesIO

# 設定日誌
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("markitdown-utils")


def reinstall_magika():
    """重新安裝 magika 套件以修復損壞的模型檔案"""
    try:
        logger.info("嘗試重新安裝 magika 套件...")
        subprocess.check_call([
            sys.executable, "-m", "pip", "uninstall", "magika", "-y"
        ])
        subprocess.check_call([
            sys.executable, "-m", "pip", "install", "magika"
        ])
        logger.info("magika 套件重新安裝完成")
        return True
    except Exception as e:
        logger.error(f"重新安裝 magika 失敗: {e}")
        return False


def convert_file_to_markdown(input_path: str,
                             use_llm: bool = False,
                             api_key: Optional[str] = None,
                             model: str = "gpt-4o") -> Tuple[bool, str,
                                                             Dict[str, Any]]:
    """
    將檔案轉換為 Markdown 格式，可選用 LLM 處理圖片
    
    Args:
        input_path (str): 輸入檔案路徑
        use_llm (bool): 是否使用 LLM 處理圖片
        api_key (str, optional): OpenAI API Key
        model (str): OpenAI 模型名稱
    
    Returns:
        Tuple[bool, str, Dict]: (是否成功, Markdown 文字, 轉換資訊)
    """
    try:
        input_path = Path(input_path).resolve()
        
        if not input_path.exists():
            logger.error(f"找不到檔案: {input_path}")
            return False, "", {"error": f"找不到檔案: {input_path}"}
            
        logger.info(f"正在轉換: {input_path}")
        
        # 檢查是否為 PPTX 檔案，如果是則使用替代方法
        if str(input_path).lower().endswith('.pptx'):
            logger.info("偵測到 PPTX 檔案，使用替代轉換方法...")
            
            # 如果啟用 LLM 且有 API Key，優先使用 Vision API 方案
            if use_llm and api_key:
                try:
                    from alternative_pptx_converter import analyze_pptx_with_vision
                    logger.info("嘗試使用 Vision API 分析 PPTX...")
                    
                    success, result_text, vision_info = analyze_pptx_with_vision(
                        str(input_path), api_key, model
                    )
                    
                    if success and result_text:
                        conversion_info = {
                            "method": "vision_api",
                            "file_name": input_path.name,
                            "file_size": input_path.stat().st_size,
                            **vision_info
                        }
                        logger.info(f"成功使用 Vision API 分析 PPTX，內容長度: {len(result_text)}")
                        return True, result_text, conversion_info
                    else:
                        logger.warning("Vision API 分析失敗，回退到 python-pptx 方法")
                        
                except ImportError:
                    logger.warning("找不到 alternative_pptx_converter 模組，使用 python-pptx 方法")
                except Exception as e:
                    logger.warning(f"Vision API 分析出錯: {e}，使用 python-pptx 方法")
            
            # 回退到原本的 python-pptx 方法
            try:
                # 嘗試使用 python-pptx 直接轉換
                from pptx import Presentation
                
                prs = Presentation(str(input_path))
                text_content = []
                slide_count = 0
                
                for slide_idx, slide in enumerate(prs.slides, 1):
                    slide_count += 1
                    text_content.append(f"\n## 投影片 {slide_idx}\n")
                    
                    slide_has_content = False
                    
                    for shape in slide.shapes:
                        # 檢查文字內容
                        if hasattr(shape, "text") and shape.text and shape.text.strip():
                            text_content.append(shape.text.strip())
                            text_content.append("")
                            slide_has_content = True
                        
                        # 檢查文字框內的段落
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                para_text = paragraph.text.strip()
                                if para_text:
                                    text_content.append(para_text)
                                    text_content.append("")
                                    slide_has_content = True
                        
                        # 檢查表格
                        if hasattr(shape, 'has_table') and shape.has_table:
                            text_content.append("\n### 表格\n")
                            table = shape.table
                            for row_idx, row in enumerate(table.rows):
                                row_text = []
                                for cell in row.cells:
                                    cell_text = cell.text.strip().replace("|", "\\|")
                                    row_text.append(cell_text)
                                text_content.append("| " + " | ".join(row_text) + " |")
                                if row_idx == 0:
                                    separator = "|" + "|".join([" --- " for _ in row.cells]) + "|"
                                    text_content.append(separator)
                            text_content.append("")
                            slide_has_content = True
                    
                    # 如果投影片沒有文字內容，檢查是否有圖片並使用 Vision API
                    if not slide_has_content:
                        image_analyzed = False
                        
                        # 如果啟用了 LLM 且有 API Key，嘗試分析圖片
                        if use_llm and api_key:
                            try:
                                # 提取投影片為圖片
                                slide_image_path = extract_slide_as_image(prs, slide_idx - 1, input_path)
                                if slide_image_path:
                                    # 使用 OpenAI Vision 分析圖片
                                    vision_result = analyze_slide_image(slide_image_path, api_key, model)
                                    if vision_result:
                                        text_content.append("### 🔍 圖片內容分析")
                                        text_content.append(vision_result)
                                        text_content.append("")
                                        image_analyzed = True
                                        slide_has_content = True
                                    
                                    # 清理臨時圖片檔案
                                    try:
                                        os.remove(slide_image_path)
                                    except:
                                        pass
                            except Exception as e:
                                logger.warning(f"分析投影片 {slide_idx} 圖片時出錯: {e}")
                        
                        # 如果沒有分析圖片或分析失敗，添加預設提示
                        if not image_analyzed:
                            text_content.append("*此投影片無文字內容或為圖片投影片*")
                            text_content.append("")
                
                result_text = "\n".join(text_content)
                
                conversion_info = {
                    "method": "python-pptx",
                    "file_name": input_path.name,
                    "file_size": input_path.stat().st_size,
                    "slide_count": slide_count,
                    "content_length": len(result_text)
                }
                
                logger.info(f"成功使用 python-pptx 轉換 {slide_count} 張投影片，內容長度: {len(result_text)}")
                
                # 如果結果為空或過短，記錄警告
                if len(result_text.strip()) < 50:
                    logger.warning(f"轉換結果可能為空或過短，內容預覽: {repr(result_text[:100])}")
                
                return True, result_text, conversion_info
                
            except ImportError:
                logger.warning("未安裝 python-pptx，嘗試使用 MarkItDown...")
            except Exception as e:
                logger.warning(f"python-pptx 轉換失敗: {e}，嘗試使用 MarkItDown...")
        
        # 建立 MarkItDown 實例
        md_kwargs = {"enable_plugins": True}
        llm_client = None
        llm_info = {}
        
        if use_llm:
            logger.info(f"嘗試啟用 LLM ({model}) 進行處理...")
            current_api_key = api_key or os.environ.get("OPENAI_API_KEY")
            if not current_api_key:
                logger.warning("未提供 OpenAI API Key，無法使用 LLM 處理圖片。")
                llm_info["status"] = "未提供 API Key"
            else:
                try:
                    llm_client = OpenAI(api_key=current_api_key)
                    # 執行一個簡單的測試呼叫來驗證金鑰
                    llm_client.models.list() 
                    logger.info("OpenAI API Key 驗證成功。")
                    md_kwargs["llm_client"] = llm_client
                    md_kwargs["llm_model"] = model
                    llm_info["status"] = "啟用成功"
                    llm_info["model"] = model
                except AuthenticationError:
                    logger.error("OpenAI API Key 無效或錯誤，無法使用 LLM。")
                    llm_info["status"] = "API Key 無效"
                except Exception as e:
                    logger.error(f"初始化 OpenAI client 時發生錯誤: {e}")
                    llm_info["status"] = f"初始化錯誤: {str(e)}"

        # 嘗試建立 MarkItDown 實例，如果失敗則嘗試修復
        max_retries = 2
        md = None
        
        for attempt in range(max_retries):
            try:
                md = MarkItDown(**md_kwargs)
                break
            except Exception as e:
                error_msg = str(e)
                logger.error(f"建立 MarkItDown 實例失敗 "
                           f"(嘗試 {attempt + 1}/{max_retries}): {error_msg}")
                
                # 檢查是否為 magika 相關錯誤
                if ("magika" in error_msg.lower() and 
                    "json" in error_msg.lower()):
                    if attempt == 0:  # 第一次嘗試修復
                        logger.info("偵測到 magika 套件問題，嘗試修復...")
                        if reinstall_magika():
                            continue
                    
                    # 如果修復失敗或是第二次嘗試，使用不依賴 magika 的方式
                    logger.warning("使用簡化模式，不使用檔案類型偵測")
                    try:
                        # 嘗試不使用 magika 的方式
                        md_kwargs_simple = {}
                        if llm_client:
                            md_kwargs_simple["llm_client"] = llm_client
                            md_kwargs_simple["llm_model"] = model
                        md = MarkItDown(**md_kwargs_simple)
                        break
                    except Exception as e2:
                        logger.error(f"簡化模式也失敗: {e2}")
                        if attempt == max_retries - 1:
                            raise e
                else:
                    if attempt == max_retries - 1:
                        raise e
        
        if md is None:
            raise Exception("無法建立 MarkItDown 實例")
        
        # 轉換檔案
        result = md.convert(str(input_path))
        
        # 準備轉換資訊
        conversion_info = {
            "llm": llm_info,
            "file_name": input_path.name,
            "file_size": input_path.stat().st_size
        }
        
        if result and result.text_content:
            # 檢查 converter_used 屬性是否存在
            if hasattr(result, 'converter_used'):
                conversion_info["converter"] = result.converter_used
            conversion_info["content_length"] = len(result.text_content)
            
            return True, result.text_content, conversion_info
        else:
            logger.error("轉換失敗，未獲得有效結果")
            return False, "", {"error": "轉換失敗，未獲得有效結果"}
            
    except Exception as e:
        logger.error(f"轉換過程中發生錯誤: {e}")
        import traceback
        error_details = traceback.format_exc()
        logger.error(error_details)
        return False, "", {"error": str(e), "details": error_details}

def convert_url_to_markdown(url: str) -> Tuple[bool, str, Dict[str, Any]]:
    """
    將 URL 轉換為 Markdown 格式
    
    Args:
        url (str): 輸入 URL
    
    Returns:
        Tuple[bool, str, Dict]: (是否成功, Markdown 文字, 轉換資訊)
    """
    try:
        logger.info(f"正在轉換 URL: {url}")
        
        # 建立 MarkItDown 實例
        md = MarkItDown(enable_plugins=True)
        
        # 轉換 URL
        result = md.convert(url)
        
        # 準備轉換資訊
        conversion_info = {
            "url": url
        }
        
        if result and result.text_content:
            # 檢查 converter_used 屬性是否存在
            if hasattr(result, 'converter_used'):
                conversion_info["converter"] = result.converter_used
            conversion_info["content_length"] = len(result.text_content)
            
            return True, result.text_content, conversion_info
        else:
            logger.error("轉換失敗，未獲得有效結果")
            return False, "", {"error": "轉換失敗，未獲得有效結果"}
            
    except Exception as e:
        logger.error(f"轉換過程中發生錯誤: {e}")
        import traceback
        error_details = traceback.format_exc()
        logger.error(error_details)
        return False, "", {"error": str(e), "details": error_details}

def extract_keywords(markdown_text: str, api_key: str, model: str = "gpt-4o-mini", count: int = 10) -> List[str]:
    """
    從 Markdown 文字中提取關鍵詞
    
    Args:
        markdown_text (str): Markdown 格式的文字
        api_key (str): OpenAI API Key
        model (str): 模型名稱
        count (int): 要提取的關鍵詞數量
    
    Returns:
        List[str]: 關鍵詞列表
    """
    try:
        client = OpenAI(api_key=api_key)
        
        response = client.chat.completions.create(
            model=model,
            messages=[
                {
                    "role": "system", 
                    "content": f"你是一個專業的關鍵詞提取工具。請從提供的文本中提取最重要的 {count} 個關鍵詞或短語，這些詞語能夠反映文本的核心主題、概念和專業術語。關鍵詞應以繁體中文提供，並按重要性排序。請僅返回關鍵詞列表，一行一個關鍵詞，不要加編號或任何其他說明。"
                },
                {
                    "role": "user",
                    "content": markdown_text
                }
            ],
            temperature=0.3
        )
        
        keywords_text = response.choices[0].message.content
        
        # 處理回應，將文字分割成列表
        keywords = [kw.strip() for kw in keywords_text.strip().split('\n') if kw.strip()]
        
        return keywords
        
    except Exception as e:
        logger.error(f"提取關鍵詞失敗: {e}")
        return []

def save_uploaded_file(uploaded_file) -> Tuple[bool, str]:
    """
    將上傳的檔案保存到臨時目錄
    
    Args:
        uploaded_file: Streamlit 上傳的檔案物件
        
    Returns:
        Tuple[bool, str]: (是否成功, 臨時檔案路徑)
    """
    try:
        # 取得檔案後綴
        file_extension = os.path.splitext(uploaded_file.name)[1]
        
        # 創建臨時檔案
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp:
            temp.write(uploaded_file.getbuffer())
            return True, temp.name
    except Exception as e:
        logger.error(f"保存上傳檔案失敗: {e}")
        return False, str(e) 

def convert_images_to_markdown(
    image_paths: List[str],
    output_file: str,
    title: str = "圖片集合",
    use_llm: bool = True,
    api_key: Optional[str] = None,
    model: str = "gpt-4o-mini"
) -> Tuple[bool, str, Dict[str, Any]]:
    """
    將多個圖片檔案轉換為單一 Markdown 檔案
    
    Args:
        image_paths (List[str]): 圖片檔案路徑列表
        output_file (str): 輸出 Markdown 檔案路徑
        title (str): Markdown 文件標題
        use_llm (bool): 是否使用 LLM 處理圖片
        api_key (Optional[str]): OpenAI API Key
        model (str): 使用的 OpenAI 模型
        
    Returns:
        Tuple[bool, str, Dict]: (是否成功, 輸出檔案路徑, 處理資訊)
    """
    try:
        logger.info(f"將 {len(image_paths)} 張圖片轉換為 Markdown")
        
        # 檢查圖片列表
        if not image_paths:
            logger.error("圖片列表為空")
            return False, "", {"error": "圖片列表為空"}
            
        # 過濾有效的圖片檔案
        valid_images = []
        for img_path in image_paths:
            if os.path.exists(img_path):
                valid_images.append(img_path)
            else:
                logger.warning(f"找不到圖片: {img_path}")
                
        if not valid_images:
            logger.error("沒有有效的圖片檔案")
            return False, "", {"error": "沒有有效的圖片檔案"}
        
        # 生成初始 Markdown 文本
        md_content = f"# {title}\n\n"
        
        # 建立 MarkItDown 實例
        md_kwargs = {"enable_plugins": True}
        llm_client = None
        llm_info = {}
        
        if use_llm:
            logger.info(f"嘗試啟用 LLM ({model}) 進行處理...")
            current_api_key = api_key or os.environ.get("OPENAI_API_KEY")
            if not current_api_key:
                logger.warning("未提供 OpenAI API Key，無法使用 LLM 處理圖片。")
                llm_info["status"] = "未提供 API Key"
                use_llm = False
            else:
                try:
                    llm_client = OpenAI(api_key=current_api_key)
                    # 執行一個簡單的測試呼叫來驗證金鑰
                    llm_client.models.list() 
                    logger.info("OpenAI API Key 驗證成功。")
                    md_kwargs["llm_client"] = llm_client
                    md_kwargs["llm_model"] = model
                    llm_info["status"] = "啟用成功"
                    llm_info["model"] = model
                except AuthenticationError:
                    logger.error("OpenAI API Key 無效或錯誤，無法使用 LLM。")
                    llm_info["status"] = "API Key 無效"
                    use_llm = False
                except Exception as e:
                    logger.error(f"初始化 OpenAI client 時發生錯誤: {e}")
                    llm_info["status"] = f"初始化錯誤: {str(e)}"
                    use_llm = False
                    
        md = MarkItDown(**md_kwargs)
        
        # 處理每個圖片
        successful_conversions = 0
        for img_path in valid_images:
            try:
                img_relpath = os.path.basename(img_path)
                logger.info(f"處理圖片: {img_relpath}")
                
                # 轉換圖片
                result = md.convert(img_path)
                
                if result and result.text_content:
                    md_content += f"## 圖片：{img_relpath}\n\n"
                    md_content += result.text_content + "\n\n"
                    successful_conversions += 1
                else:
                    logger.warning(f"無法轉換圖片: {img_relpath}")
                    # 添加簡單的圖片標記
                    md_content += f"## 圖片：{img_relpath}\n\n"
                    md_content += f"![{img_relpath}]({img_path})\n\n"
            except Exception as e:
                logger.warning(f"處理圖片 {img_path} 時出錯: {e}")
                # 添加簡單的圖片標記
                md_content += f"## 圖片：{img_relpath}\n\n"
                md_content += f"![{img_relpath}]({img_path})\n\n"
                
        # 寫入輸出檔案
        try:
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(md_content)
            logger.info(f"已將 {successful_conversions} 張圖片轉換並寫入 {output_file}")
            
            # 嘗試使用 image_analyzer 增強圖片描述
            if use_llm and os.path.exists(output_file):
                try:
                    from image_analyzer import enhance_markdown_with_image_analysis
                    
                    # 讀取剛生成的 Markdown 檔案
                    with open(output_file, 'r', encoding='utf-8') as f:
                        original_md = f.read()
                    
                    # 增強 Markdown 內容
                    enhanced_md, stats = enhance_markdown_with_image_analysis(
                        markdown_text=original_md,
                        base_dir=os.path.dirname(output_file),
                        api_key=current_api_key,
                        model=model
                    )
                    
                    # 寫入增強後的內容
                    with open(output_file, 'w', encoding='utf-8') as f:
                        f.write(enhanced_md)
                    
                    logger.info(f"已增強 {stats['images_processed']} 張圖片的描述")
                    
                except ImportError:
                    logger.warning("找不到 image_analyzer 模組，無法增強圖片描述")
                except Exception as e:
                    logger.warning(f"增強圖片描述時出錯: {e}")
            
            result_info = {
                "success": True,
                "total_images": len(valid_images),
                "converted_images": successful_conversions,
                "output_file": output_file,
                "llm": llm_info
            }
            
            return True, output_file, result_info
            
        except Exception as e:
            logger.error(f"寫入輸出檔案時發生錯誤: {e}")
            return False, "", {"error": f"寫入輸出檔案時發生錯誤: {e}"}
        
    except Exception as e:
        logger.error(f"轉換圖片集合時發生錯誤: {e}")
        import traceback
        error_details = traceback.format_exc()
        logger.error(error_details)
        return False, "", {"error": str(e), "details": error_details}


def extract_slide_as_image(presentation, slide_index: int, input_path: str) -> Optional[str]:
    """
    將 PowerPoint 投影片轉換為圖片
    
    Args:
        presentation: python-pptx Presentation 物件
        slide_index: 投影片索引 (0-based)
        input_path: 原始 PPTX 檔案路徑
        
    Returns:
        Optional[str]: 臨時圖片檔案路徑，失敗時返回 None
    """
    try:
        # 這裡需要使用其他方法來將投影片轉為圖片
        # 由於 python-pptx 不直接支援轉圖片，我們可以使用其他方案
        
        # 方案 1: 嘗試使用 PIL 和 python-pptx 的形狀資訊
        # 但這個方法有限制，更好的方案是使用外部工具
        
        # 暫時返回 None，表示無法提取圖片
        # 在實際使用中，可以整合 LibreOffice 或其他工具
        logger.warning(f"投影片 {slide_index + 1} 圖片提取功能尚未完全實現")
        return None
        
    except Exception as e:
        logger.error(f"提取投影片 {slide_index + 1} 為圖片時出錯: {e}")
        return None


def analyze_slide_image(image_path: str, api_key: str, model: str = "gpt-4o") -> Optional[str]:
    """
    使用 OpenAI Vision API 分析投影片圖片
    
    Args:
        image_path: 圖片檔案路徑
        api_key: OpenAI API Key
        model: 使用的模型名稱
        
    Returns:
        Optional[str]: 分析結果文字，失敗時返回 None
    """
    try:
        client = OpenAI(api_key=api_key)
        
        # 讀取並編碼圖片
        with open(image_path, "rb") as image_file:
            base64_image = base64.b64encode(image_file.read()).decode('utf-8')
        
        response = client.chat.completions.create(
            model=model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "請仔細分析這張投影片圖片的內容，包括文字、圖表、圖像等元素，並用繁體中文詳細描述。請特別注意：1) 提取所有可見的文字內容 2) 描述圖表、圖像的類型和內容 3) 解釋投影片的主要訊息和重點"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=1000,
            temperature=0.3
        )
        
        result = response.choices[0].message.content
        logger.info(f"成功分析投影片圖片，結果長度: {len(result) if result else 0}")
        return result
        
    except Exception as e:
        logger.error(f"使用 Vision API 分析投影片圖片時出錯: {e}")
        return None 